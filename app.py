"""
General Chatbot — Multi-Provider Edition
==========================================
A single-file Streamlit chatbot that supports:
  - Local models via Ollama
  - Google Gemini
  - OpenAI (GPT models)
  - Anthropic (Claude models)
  - PDF / DOCX upload so you can chat about a document's content

Run with:
    streamlit run app.py
"""

import streamlit as st
import requests
import json
import io
import re
import sqlite3
import uuid
import base64
import subprocess
import tempfile
from pathlib import Path
from datetime import datetime
from bs4 import BeautifulSoup

# ---- Optional imports for file parsing (only needed if user uploads files) ----
try:
    import fitz  # PyMuPDF
    PDF_SUPPORT = True
except ImportError:
    PDF_SUPPORT = False

try:
    import docx  # python-docx
    DOCX_SUPPORT = True
except ImportError:
    DOCX_SUPPORT = False

try:
    from pptx import Presentation as PptxPresentation
    PPTX_SUPPORT = True
except ImportError:
    PPTX_SUPPORT = False

try:
    import openpyxl
    XLSX_SUPPORT = True
except ImportError:
    XLSX_SUPPORT = False

try:
    from pdf2image import convert_from_path
    PDF2IMAGE_SUPPORT = True
except ImportError:
    PDF2IMAGE_SUPPORT = False

# LibreOffice is a system binary, not a pip package — check both common command names.
import shutil as _shutil
LIBREOFFICE_BIN = _shutil.which("soffice") or _shutil.which("libreoffice")
LIBREOFFICE_SUPPORT = LIBREOFFICE_BIN is not None


# =============================================================================
# PAGE CONFIG
# =============================================================================
st.set_page_config(
    page_title="General Chatbot",
    page_icon="💬",
    layout="wide",
)


# =============================================================================
# CONSTANTS
# =============================================================================
OLLAMA_BASE_URL = "http://localhost:11434"

GEMINI_MODELS = ["gemini-2.5-flash", "gemini-2.5-pro"]
OPENAI_MODELS = ["gpt-4o", "gpt-4o-mini", "gpt-4.1"]
ANTHROPIC_MODELS = ["claude-sonnet-4-6", "claude-opus-4-7", "claude-haiku-4-5-20251001"]
GROQ_MODELS = ["openai/gpt-oss-120b", "openai/gpt-oss-20b", "meta-llama/llama-4-scout-17b-16e-instruct"]

# Groq vision support is currently limited to specific multimodal models (e.g. Llama 4 Scout).
GROQ_VISION_MODEL_HINTS = ["scout", "maverick", "vision"]

MAX_DOC_CHARS = 15000  # cap injected document text to keep prompts manageable
MAX_DOCS_ACTIVE = 10   # safety cap on number of simultaneously active documents

# ---- Persistent chat history ----
DB_PATH = str(Path(__file__).parent / "chat_history.db")

# ---- Vision support ----
# Gemini, OpenAI, and Anthropic all support image input on their standard models.
# Ollama only supports images on specific multimodal models — text-only models
# (llama3.2, mistral, etc.) will silently ignore or error on image input, so we
# warn the user instead of pretending it works.
OLLAMA_VISION_MODEL_HINTS = ["llava", "vision", "bakllava", "moondream"]

# ---- Supported file categories ----
TEXT_FILE_EXTENSIONS = {".txt", ".md", ".py"}
IMAGE_FILE_EXTENSIONS = {".png", ".jpg", ".jpeg"}
SLIDE_SHEET_EXTENSIONS = {".pptx", ".xlsx"}
ALL_SUPPORTED_EXTENSIONS = (
    {".pdf", ".docx"} | TEXT_FILE_EXTENSIONS | IMAGE_FILE_EXTENSIONS | SLIDE_SHEET_EXTENSIONS
)

# ---- Web search auto-detect settings ----
# If the user's message contains any of these signals, we attempt a live web search
# before answering, since it suggests they want current/real-time information.
SEARCH_TRIGGER_KEYWORDS = [
    "latest", "current", "currently", "today", "right now", "as of now",
    "recent", "recently", "this week", "this month", "this year",
    "news", "breaking", "update", "updated",
    "price of", "stock price", "exchange rate", "weather",
    "who is the current", "who is the present", "what is the current",
    "score", "result of", "live",
    "2026", "2027",  # catches "in 2026", "latest 2026", etc.
]
SEARCH_RESULT_COUNT = 4  # how many DuckDuckGo results to pull in

# ---- Visual theme: provider identity (color-coded chips, used in sidebar + bubbles) ----
PROVIDER_STYLE = {
    "Ollama (local)": {"color": "#F0997B", "bg": "#2A1F17", "border": "#D85A30", "icon": "🦙", "label": "ollama"},
    "Gemini":         {"color": "#85B7EB", "bg": "#16202B", "border": "#378ADD", "icon": "✨", "label": "gemini"},
    "OpenAI":         {"color": "#97C459", "bg": "#1A1F17", "border": "#639922", "icon": "◎", "label": "openai"},
    "Anthropic":      {"color": "#ED93B1", "bg": "#1F1A1D", "border": "#D4537E", "icon": "✦", "label": "anthropic"},
    "Groq":           {"color": "#F0B27A", "bg": "#241C14", "border": "#E8772E", "icon": "⚡", "label": "groq"},
}
USER_AVATAR_BG = "#3D3729"
USER_AVATAR_COLOR = "#C9C2B0"


# =============================================================================
# FILE EXTRACTION HELPERS
# =============================================================================
def extract_pdf_text(file_bytes: bytes) -> str:
    """Extract plain text from a PDF using PyMuPDF."""
    if not PDF_SUPPORT:
        return "[Error: PyMuPDF not installed. Run: pip install pymupdf --break-system-packages]"
    try:
        doc = fitz.open(stream=file_bytes, filetype="pdf")
        text_parts = []
        for page_num, page in enumerate(doc, start=1):
            page_text = page.get_text()
            if page_text.strip():
                text_parts.append(f"--- Page {page_num} ---\n{page_text}")
        doc.close()
        return "\n\n".join(text_parts) if text_parts else "[No extractable text found in PDF]"
    except Exception as e:
        return f"[Error reading PDF: {e}]"


def extract_docx_text(file_bytes: bytes) -> str:
    """Extract plain text from a DOCX using python-docx."""
    if not DOCX_SUPPORT:
        return "[Error: python-docx not installed. Run: pip install python-docx --break-system-packages]"
    try:
        document = docx.Document(io.BytesIO(file_bytes))
        paragraphs = [p.text for p in document.paragraphs if p.text.strip()]
        # Also pull text out of tables, since they're common in real documents
        for table in document.tables:
            for row in table.rows:
                row_text = " | ".join(cell.text.strip() for cell in row.cells if cell.text.strip())
                if row_text:
                    paragraphs.append(row_text)
        return "\n".join(paragraphs) if paragraphs else "[No extractable text found in DOCX]"
    except Exception as e:
        return f"[Error reading DOCX: {e}]"


def needs_web_search(text: str) -> bool:
    """
    Auto-detect whether a user message likely needs current/real-time information.
    This is a simple keyword heuristic, not a model call, so it's instant and free.
    """
    lowered = text.lower()
    return any(keyword in lowered for keyword in SEARCH_TRIGGER_KEYWORDS)


def web_search_duckduckgo(query: str, max_results: int = SEARCH_RESULT_COUNT) -> str:
    """
    Attempt a live web search via DuckDuckGo's HTML endpoint (no API key required).

    IMPORTANT — honest limitation:
    DuckDuckGo does not offer an official free search API. This function scrapes
    their HTML results page, which:
      - Is against DuckDuckGo's Terms of Service for automated use
      - Has no published rate limit, but community reports suggest it can start
        blocking/rate-limiting (HTTP 202/403/429) after roughly 20-30 requests/minute
        from a single IP
      - Can break entirely if DuckDuckGo changes their page's HTML structure

    This function is built to FAIL GRACEFULLY: if anything goes wrong, it returns
    None instead of crashing, so the chatbot still answers (just without fresh
    context) rather than showing an error to the user.

    For a production app, swap this for a real search API (e.g. Tavily's free
    tier) — see README for notes on how to do that.
    """
    try:
        headers = {
            "User-Agent": (
                "Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 "
                "(KHTML, like Gecko) Chrome/124.0.0.0 Safari/537.36"
            )
        }
        response = requests.post(
            "https://html.duckduckgo.com/html/",
            data={"q": query},
            headers=headers,
            timeout=10,
        )

        # DuckDuckGo signals rate-limiting / soft-blocking with these statuses
        if response.status_code in (202, 403, 429):
            return None

        response.raise_for_status()

        soup = BeautifulSoup(response.text, "html.parser")
        results = soup.find_all("div", class_="result__body")

        if not results:
            # Page loaded but markup didn't match what we expect —
            # DuckDuckGo likely changed their HTML structure.
            return None

        snippets = []
        for result in results[:max_results]:
            title_tag = result.find("a", class_="result__a")
            snippet_tag = result.find("a", class_="result__snippet")

            title = title_tag.get_text(strip=True) if title_tag else ""
            snippet = snippet_tag.get_text(strip=True) if snippet_tag else ""

            if title or snippet:
                snippets.append(f"- {title}: {snippet}")

        if not snippets:
            return None

        return "\n".join(snippets)

    except requests.exceptions.RequestException:
        # Network error, timeout, connection refused, etc.
        return None
    except Exception:
        # Catch-all so a parsing quirk never crashes the chat
        return None

def render_office_file_to_images(file_bytes: bytes, suffix: str, dpi: int = 130) -> list:
    """
    Convert a PPTX or XLSX file to a list of base64-encoded PNG images, one per
    slide/page, using LibreOffice (headless) to produce a PDF, then pdf2image to
    rasterize each page.

    Returns [] if LibreOffice or pdf2image aren't available, or if conversion
    fails for any reason — callers should treat that as "no visual available"
    and fall back to text-only context, not as a crash.
    """
    if not (LIBREOFFICE_SUPPORT and PDF2IMAGE_SUPPORT):
        return []

    try:
        with tempfile.TemporaryDirectory() as tmpdir:
            tmp_input = Path(tmpdir) / f"input{suffix}"
            tmp_input.write_bytes(file_bytes)

            # LibreOffice writes <basename>.pdf into --outdir using the input's basename
            result = subprocess.run(
                [LIBREOFFICE_BIN, "--headless", "--convert-to", "pdf", "--outdir", tmpdir, str(tmp_input)],
                capture_output=True,
                timeout=90,
            )
            pdf_path = Path(tmpdir) / "input.pdf"
            if not pdf_path.exists():
                return []

            images = convert_from_path(str(pdf_path), dpi=dpi)
            encoded_images = []
            for img in images:
                buf = io.BytesIO()
                img.save(buf, format="PNG")
                encoded_images.append(base64.b64encode(buf.getvalue()).decode("utf-8"))
            return encoded_images

    except (subprocess.TimeoutExpired, Exception):
        # Broad catch is intentional here: any failure in this optional visual
        # pipeline should degrade to "no images", never break the chat.
        return []


def extract_pptx_text(file_bytes: bytes) -> str:
    """Extract plain text (titles + body text) from a PPTX using python-pptx."""
    if not PPTX_SUPPORT:
        return "[Error: python-pptx not installed. Run: pip install python-pptx --break-system-packages]"
    try:
        prs = PptxPresentation(io.BytesIO(file_bytes))
        slide_texts = []
        for i, slide in enumerate(prs.slides, start=1):
            parts = []
            for shape in slide.shapes:
                if shape.has_text_frame and shape.text_frame.text.strip():
                    parts.append(shape.text_frame.text.strip())
            if parts:
                slide_texts.append(f"--- Slide {i} ---\n" + "\n".join(parts))
        return "\n\n".join(slide_texts) if slide_texts else "[No extractable text found in PPTX]"
    except Exception as e:
        return f"[Error reading PPTX: {e}]"


def extract_xlsx_text(file_bytes: bytes) -> str:
    """Extract cell values from every sheet in an XLSX using openpyxl."""
    if not XLSX_SUPPORT:
        return "[Error: openpyxl not installed. Run: pip install openpyxl --break-system-packages]"
    try:
        wb = openpyxl.load_workbook(io.BytesIO(file_bytes), data_only=True)
        sheet_texts = []
        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            rows = []
            for row in ws.iter_rows(values_only=True):
                if any(cell is not None for cell in row):
                    row_str = " | ".join(str(cell) if cell is not None else "" for cell in row)
                    rows.append(row_str)
            if rows:
                sheet_texts.append(f"--- Sheet: {sheet_name} ---\n" + "\n".join(rows))
        return "\n\n".join(sheet_texts) if sheet_texts else "[No data found in XLSX]"
    except Exception as e:
        return f"[Error reading XLSX: {e}]"


def process_uploaded_file(uploaded_file) -> dict:
    """
    Route an uploaded file to the right extractor(s) based on its extension.

    Returns a dict: {"name", "ext", "text", "images", "image_render_attempted",
    "image_render_failed"} where:
      - "text" holds extracted plain text (may be empty for pure images)
      - "images" is a list of base64-encoded PNGs (populated for image files,
        and for PPTX/XLSX when LibreOffice rendering succeeds)
      - "image_render_attempted"/"image_render_failed" let the UI explain when
        a PPTX/XLSX was uploaded but visual rendering wasn't available
    """
    file_bytes = uploaded_file.read()
    name = uploaded_file.name
    ext = Path(name).suffix.lower()

    result = {
        "name": name,
        "ext": ext,
        "text": "",
        "images": [],
        "image_render_attempted": False,
        "image_render_failed": False,
    }

    if ext == ".pdf":
        result["text"] = extract_pdf_text(file_bytes)

    elif ext == ".docx":
        result["text"] = extract_docx_text(file_bytes)

    elif ext in TEXT_FILE_EXTENSIONS:
        text = file_bytes.decode("utf-8", errors="ignore")
        if ext == ".py":
            text = f"```python\n{text}\n```"
        result["text"] = text

    elif ext in IMAGE_FILE_EXTENSIONS:
        result["images"] = [base64.b64encode(file_bytes).decode("utf-8")]
        result["text"] = ""  # no separate text; the image itself is the content

    elif ext == ".pptx":
        result["text"] = extract_pptx_text(file_bytes)
        result["image_render_attempted"] = True
        result["images"] = render_office_file_to_images(file_bytes, ".pptx")
        result["image_render_failed"] = not result["images"]

    elif ext == ".xlsx":
        result["text"] = extract_xlsx_text(file_bytes)
        result["image_render_attempted"] = True
        result["images"] = render_office_file_to_images(file_bytes, ".xlsx")
        result["image_render_failed"] = not result["images"]

    else:
        result["text"] = f"[Unsupported file type: {ext}]"

    if len(result["text"]) > MAX_DOC_CHARS:
        result["text"] = result["text"][:MAX_DOC_CHARS] + "\n\n[... content truncated due to length ...]"

    return result


# =============================================================================
# PROVIDER CALL FUNCTIONS
# Each function takes (messages, model, api_key) and returns a string reply.
# `messages` is a list of dicts: {"role": "user"/"assistant", "content": "...",
# "images": [base64_png, ...]}  -- "images" is optional and only ever present
# on user messages that had a file/image attached.
# =============================================================================

def call_ollama(messages, model, api_key=None):
    """
    Call a locally running Ollama model. No API key needed.
    Ollama's /api/chat accepts an "images" field (list of base64 strings,
    no data URI prefix) on a message — but only multimodal models actually
    use it. We still send it if present; non-vision models will just ignore
    or error on it, which is why the UI warns the user beforehand.
    """
    try:
        ollama_messages = []
        for m in messages:
            entry = {"role": m["role"], "content": m["content"]}
            if m.get("images"):
                entry["images"] = m["images"]
            ollama_messages.append(entry)

        response = requests.post(
            f"{OLLAMA_BASE_URL}/api/chat",
            json={"model": model, "messages": ollama_messages, "stream": False},
            timeout=180,
        )
        response.raise_for_status()
        data = response.json()
        return data.get("message", {}).get("content", "[No content returned by Ollama]")
    except requests.exceptions.ConnectionError:
        return (
            "⚠️ Could not connect to Ollama at http://localhost:11434. "
            "Make sure Ollama is running (`ollama serve`) and the model is pulled "
            f"(`ollama pull {model}`)."
        )
    except Exception as e:
        return f"⚠️ Ollama error: {e}"


def call_openai(messages, model, api_key):
    """
    Call OpenAI's Chat Completions API.
    For image input, OpenAI expects content as a list of parts:
    [{"type": "text", "text": "..."}, {"type": "image_url", "image_url": {"url": "data:image/png;base64,..."}}]
    """
    if not api_key:
        return "⚠️ Please enter your OpenAI API key in the sidebar."
    try:
        openai_messages = []
        for m in messages:
            if m.get("images"):
                parts = [{"type": "text", "text": m["content"]}]
                for img_b64 in m["images"]:
                    parts.append({
                        "type": "image_url",
                        "image_url": {"url": f"data:image/png;base64,{img_b64}"},
                    })
                openai_messages.append({"role": m["role"], "content": parts})
            else:
                openai_messages.append({"role": m["role"], "content": m["content"]})

        response = requests.post(
            "https://api.openai.com/v1/chat/completions",
            headers={
                "Authorization": f"Bearer {api_key}",
                "Content-Type": "application/json",
            },
            json={"model": model, "messages": openai_messages, "max_tokens": 2000},
            timeout=60,
        )
        response.raise_for_status()
        data = response.json()
        return data["choices"][0]["message"]["content"]
    except requests.exceptions.HTTPError as e:
        return f"⚠️ OpenAI API error: {e.response.status_code} - {e.response.text}"
    except Exception as e:
        return f"⚠️ OpenAI error: {e}"


def call_groq(messages, model, api_key):
    """
    Call Groq's Chat Completions API (OpenAI-compatible format).
    For image input, Groq expects content as a list of parts, same as OpenAI:
    [{"type": "text", "text": "..."}, {"type": "image_url", "image_url": {"url": "data:image/png;base64,..."}}]
    """
    if not api_key:
        return "⚠️ Please enter your Groq API key in the sidebar."
    try:
        groq_messages = []
        for m in messages:
            if m.get("images"):
                parts = [{"type": "text", "text": m["content"]}]
                for img_b64 in m["images"]:
                    parts.append({
                        "type": "image_url",
                        "image_url": {"url": f"data:image/png;base64,{img_b64}"},
                    })
                groq_messages.append({"role": m["role"], "content": parts})
            else:
                groq_messages.append({"role": m["role"], "content": m["content"]})

        response = requests.post(
            "https://api.groq.com/openai/v1/chat/completions",
            headers={
                "Authorization": f"Bearer {api_key}",
                "Content-Type": "application/json",
            },
            json={"model": model, "messages": groq_messages, "max_tokens": 2000},
            timeout=60,
        )
        response.raise_for_status()
        data = response.json()
        return data["choices"][0]["message"]["content"]
    except requests.exceptions.HTTPError as e:
        return f"⚠️ Groq API error: {e.response.status_code} - {e.response.text}"
    except Exception as e:
        return f"⚠️ Groq error: {e}"


def call_anthropic(messages, model, api_key):
    """
    Call Anthropic's Messages API.
    For image input, Anthropic expects content as a list of blocks:
    [{"type": "image", "source": {"type": "base64", "media_type": "image/png", "data": "..."}},
     {"type": "text", "text": "..."}]
    """
    if not api_key:
        return "⚠️ Please enter your Anthropic API key in the sidebar."
    try:
        anthropic_messages = []
        for m in messages:
            if m.get("images"):
                blocks = []
                for img_b64 in m["images"]:
                    blocks.append({
                        "type": "image",
                        "source": {"type": "base64", "media_type": "image/png", "data": img_b64},
                    })
                blocks.append({"type": "text", "text": m["content"]})
                anthropic_messages.append({"role": m["role"], "content": blocks})
            else:
                anthropic_messages.append({"role": m["role"], "content": m["content"]})

        response = requests.post(
            "https://api.anthropic.com/v1/messages",
            headers={
                "x-api-key": api_key,
                "anthropic-version": "2023-06-01",
                "Content-Type": "application/json",
            },
            json={
                "model": model,
                "max_tokens": 2000,
                "messages": anthropic_messages,
            },
            timeout=60,
        )
        response.raise_for_status()
        data = response.json()
        return data["content"][0]["text"]
    except requests.exceptions.HTTPError as e:
        return f"⚠️ Anthropic API error: {e.response.status_code} - {e.response.text}"
    except Exception as e:
        return f"⚠️ Anthropic error: {e}"


def call_gemini(messages, model, api_key):
    """
    Call Google's Gemini API. Gemini uses 'contents' with role user/model.
    For image input, parts include an inline_data block (REST API uses
    snake_case field names — camelCase like 'inlineData' only works in the
    official SDKs, which translate it internally; raw HTTP needs snake_case):
    {"inline_data": {"mime_type": "image/png", "data": "..."}}
    """
    if not api_key:
        return "⚠️ Please enter your Gemini API key in the sidebar."
    try:
        contents = []
        for m in messages:
            role = "model" if m["role"] == "assistant" else "user"
            parts = [{"text": m["content"]}]
            for img_b64 in m.get("images", []):
                parts.append({"inline_data": {"mime_type": "image/png", "data": img_b64}})
            contents.append({"role": role, "parts": parts})

        url = (
            f"https://generativelanguage.googleapis.com/v1beta/models/"
            f"{model}:generateContent?key={api_key}"
        )
        response = requests.post(
            url,
            headers={"Content-Type": "application/json"},
            json={"contents": contents},
            timeout=60,
        )
        response.raise_for_status()
        data = response.json()
        return data["candidates"][0]["content"]["parts"][0]["text"]
    except requests.exceptions.HTTPError as e:
        return f"⚠️ Gemini API error: {e.response.status_code} - {e.response.text}"
    except (KeyError, IndexError):
        return "⚠️ Gemini returned an unexpected response (possibly blocked by safety filters)."
    except Exception as e:
        return f"⚠️ Gemini error: {e}"


def get_ollama_models():
    """Fetch list of locally available Ollama models. Returns [] if Ollama is unreachable."""
    try:
        response = requests.get(f"{OLLAMA_BASE_URL}/api/tags", timeout=5)
        response.raise_for_status()
        data = response.json()
        return [m["name"] for m in data.get("models", [])]
    except Exception:
        return []


def is_ollama_vision_model(model_name: str) -> bool:
    """Heuristic check for whether an Ollama model name suggests vision support."""
    lowered = model_name.lower()
    return any(hint in lowered for hint in OLLAMA_VISION_MODEL_HINTS)


def is_groq_vision_model(model_name: str) -> bool:
    """Heuristic check for whether a Groq model name suggests vision support."""
    lowered = model_name.lower()
    return any(hint in lowered for hint in GROQ_VISION_MODEL_HINTS)


# =============================================================================
# PERSISTENT CHAT HISTORY (SQLite)
# =============================================================================
def init_db():
    """Create the chats table if it doesn't already exist. Safe to call every run."""
    conn = sqlite3.connect(DB_PATH)
    conn.execute(
        """
        CREATE TABLE IF NOT EXISTS chats (
            id TEXT PRIMARY KEY,
            title TEXT NOT NULL,
            created_at TEXT NOT NULL,
            updated_at TEXT NOT NULL,
            messages_json TEXT NOT NULL
        )
        """
    )
    conn.commit()
    conn.close()


def save_chat_to_db(chat_id: str, title: str, messages: list):
    """Insert a new chat or update an existing one (matched by id)."""
    if not messages:
        return  # don't save empty chats
    conn = sqlite3.connect(DB_PATH)
    now = datetime.now().isoformat()
    messages_json = json.dumps(messages)
    existing = conn.execute("SELECT id FROM chats WHERE id = ?", (chat_id,)).fetchone()
    if existing:
        conn.execute(
            "UPDATE chats SET title = ?, updated_at = ?, messages_json = ? WHERE id = ?",
            (title, now, messages_json, chat_id),
        )
    else:
        conn.execute(
            "INSERT INTO chats (id, title, created_at, updated_at, messages_json) VALUES (?, ?, ?, ?, ?)",
            (chat_id, title, now, now, messages_json),
        )
    conn.commit()
    conn.close()


def list_saved_chats():
    """Return [(id, title, updated_at), ...] ordered most-recently-updated first."""
    conn = sqlite3.connect(DB_PATH)
    rows = conn.execute(
        "SELECT id, title, updated_at FROM chats ORDER BY updated_at DESC"
    ).fetchall()
    conn.close()
    return rows


def load_chat_from_db(chat_id: str):
    """Return the messages list for a saved chat, or None if not found."""
    conn = sqlite3.connect(DB_PATH)
    row = conn.execute("SELECT messages_json FROM chats WHERE id = ?", (chat_id,)).fetchone()
    conn.close()
    return json.loads(row[0]) if row else None


def delete_chat_from_db(chat_id: str):
    conn = sqlite3.connect(DB_PATH)
    conn.execute("DELETE FROM chats WHERE id = ?", (chat_id,))
    conn.commit()
    conn.close()


def generate_chat_title(first_message: str) -> str:
    """Turn the first user message into a short chat title, ChatGPT-style."""
    title = first_message.strip().replace("\n", " ")
    if len(title) > 50:
        title = title[:50].rsplit(" ", 1)[0] + "..."
    return title or "New chat"


def escape_html(text: str) -> str:
    """
    Escape HTML special characters so user-typed or model-generated text can
    never be interpreted as raw HTML/JS inside our custom bubble markup.
    This is essential: without it, a message containing '<script>' or stray
    '<'/'>' characters could break the layout or, worse, inject markup.
    """
    return (
        text.replace("&", "&amp;")
        .replace("<", "&lt;")
        .replace(">", "&gt;")
        .replace('"', "&quot;")
    )


def format_message_html(text: str) -> str:
    """
    Escape the text for safety, then apply a few lightweight Markdown-style
    conversions (bold, inline code, line breaks) so replies still look
    readable inside our custom HTML bubbles, which bypass st.markdown's
    normal renderer.
    """
    safe = escape_html(text)
    # Inline code: `code` -> <code>code</code>
    safe = re.sub(r"`([^`\n]+)`", r"<code>\1</code>", safe)
    # Bold: **text** -> <strong>text</strong>
    safe = re.sub(r"\*\*([^*\n]+)\*\*", r"<strong>\1</strong>", safe)
    # Preserve line breaks
    safe = safe.replace("\n", "<br>")
    return safe


def inject_theme_css():
    """
    Inject the custom dark + warm-coral theme.
    Streamlit's default chat UI is overridden here so we can render custom
    HTML bubbles with avatars and timestamps instead of the built-in widget.
    """
    st.markdown(
        """
        <style>
        /* ---- App background & base text ---- */
        .stApp {
            background-color: #14120F;
        }
        section[data-testid="stSidebar"] {
            background-color: #1A1813;
            border-right: 1px solid #2A2722;
        }
        h1, h2, h3, p, span, label, .stMarkdown {
            color: #EDEAE2;
        }
        .stCaption, [data-testid="stCaptionContainer"] {
            color: #7A7363 !important;
        }

        /* ---- Inputs, selects, buttons ---- */
        .stTextInput input, .stSelectbox div[data-baseweb="select"] > div {
            background-color: #1F1C17 !important;
            border: 1px solid #3D3729 !important;
            color: #EDEAE2 !important;
        }
        .stButton button {
            background-color: #1F1C17;
            border: 1px solid #3D3729;
            color: #EDEAE2;
        }
        .stButton button:hover {
            border-color: #D85A30;
            color: #D85A30;
        }
        .stFileUploader section {
            background-color: #1A1813;
            border: 1px dashed #3D3729;
        }

        /* ---- Chat input bar ---- */
        .stChatInput textarea, [data-testid="stChatInput"] textarea {
            background-color: #1A1813 !important;
            color: #EDEAE2 !important;
            border: 1px solid #3D3729 !important;
            border-radius: 20px !important;
        }
        [data-testid="stChatInput"] {
            background-color: #14120F !important;
        }

        /* ---- Hide default streamlit chrome we don't need ---- */
        [data-testid="stChatMessageAvatarUser"], [data-testid="stChatMessageAvatarAssistant"] {
            display: none;
        }

        /* ---- Custom chat bubble layout ---- */
        .chat-row {
            display: flex;
            gap: 10px;
            margin-bottom: 16px;
            align-items: flex-start;
        }
        .chat-row.user {
            justify-content: flex-end;
        }
        .avatar {
            width: 30px;
            height: 30px;
            border-radius: 50%;
            display: flex;
            align-items: center;
            justify-content: center;
            font-size: 13px;
            flex-shrink: 0;
            font-family: -apple-system, sans-serif;
        }
        .bubble {
            padding: 11px 16px;
            max-width: 72%;
            font-size: 14.5px;
            line-height: 1.6;
        }
        .bubble.user {
            background-color: #23201A;
            border-radius: 16px 16px 3px 16px;
            color: #EDEAE2;
        }
        .bubble.assistant {
            background-color: #1A1813;
            border-radius: 3px 16px 16px 16px;
            color: #D6D1C4;
            border: 1px solid #242019;
        }
        .bubble-meta {
            display: flex;
            align-items: center;
            gap: 6px;
            margin-bottom: 7px;
            font-size: 11px;
            color: #7A7363;
            font-family: -apple-system, sans-serif;
        }
        .bubble-timestamp {
            display: block;
            margin-top: 6px;
            font-size: 10px;
            color: #5A5547;
            font-family: -apple-system, sans-serif;
        }
        .search-indicator {
            font-size: 11px;
            color: #1D9E75;
        }
        </style>
        """,
        unsafe_allow_html=True,
    )


def render_bubble(role: str, content: str, timestamp: str, provider_info: dict = None, search_used: bool = None):
    """
    Render one chat message as a custom HTML bubble with avatar + timestamp,
    replacing Streamlit's default st.chat_message styling.
    """
    if role == "user":
        st.markdown(
            f'<div class="chat-row user">'
            f'<div class="bubble user">'
            f'<div>{format_message_html(content)}</div>'
            f'<span class="bubble-timestamp">{timestamp}</span>'
            f'</div>'
            f'<div class="avatar" style="background-color:{USER_AVATAR_BG}; color:{USER_AVATAR_COLOR};">you</div>'
            f'</div>',
            unsafe_allow_html=True,
        )
    else:
        style = provider_info or PROVIDER_STYLE["Anthropic"]
        meta_html = ""
        if search_used is True:
            meta_html = '<div class="bubble-meta"><span class="search-indicator">●</span> live search used</div>'
        elif search_used is False:
            meta_html = '<div class="bubble-meta">⚠ live search unavailable — general knowledge used</div>'

        st.markdown(
            f'<div class="chat-row assistant">'
            f'<div class="avatar" style="background-color:{style["bg"]}; color:{style["color"]}; border:1px solid {style["border"]};">{style["icon"]}</div>'
            f'<div class="bubble assistant">'
            f'{meta_html}'
            f'<div>{format_message_html(content)}</div>'
            f'<span class="bubble-timestamp">{timestamp}</span>'
            f'</div>'
            f'</div>',
            unsafe_allow_html=True,
        )


# =============================================================================
# THEME
# =============================================================================
inject_theme_css()
init_db()


# =============================================================================
# SESSION STATE INIT
# =============================================================================
if "messages" not in st.session_state:
    st.session_state.messages = []  # list of {role, content, timestamp, provider, search_used, images}

if "documents" not in st.session_state:
    st.session_state.documents = []  # list of dicts from process_uploaded_file()

if "processed_file_names" not in st.session_state:
    st.session_state.processed_file_names = set()  # avoid re-processing same upload on rerun

if "current_chat_id" not in st.session_state:
    st.session_state.current_chat_id = str(uuid.uuid4())


# =============================================================================
# SIDEBAR — NEW CHAT / RECENT / PROVIDER / MODEL / API KEY / FILE UPLOAD
# =============================================================================
with st.sidebar:
    st.markdown("### ✦ General chatbot")
    st.caption("Multi-provider AI chat")

    # ---- New chat ----
    if st.button("＋  New chat", use_container_width=True):
        # Save the current chat before starting a new one, so nothing is lost
        if st.session_state.messages:
            title = generate_chat_title(st.session_state.messages[0]["content"])
            save_chat_to_db(st.session_state.current_chat_id, title, st.session_state.messages)
        st.session_state.messages = []
        st.session_state.documents = []
        st.session_state.processed_file_names = set()
        st.session_state.current_chat_id = str(uuid.uuid4())
        st.rerun()

    # ---- Recent chats ----
    with st.expander("🕘  Recent", expanded=False):
        saved_chats = list_saved_chats()
        if not saved_chats:
            st.caption("No saved chats yet — they appear here once you start one.")
        for chat_id, title, updated_at in saved_chats[:30]:
            is_current = chat_id == st.session_state.current_chat_id
            col1, col2 = st.columns([5, 1])
            with col1:
                label = f"**{title}**" if is_current else title
                if st.button(label, key=f"open_{chat_id}", use_container_width=True):
                    # Save current chat first so switching away doesn't lose it
                    if st.session_state.messages and not is_current:
                        cur_title = generate_chat_title(st.session_state.messages[0]["content"])
                        save_chat_to_db(st.session_state.current_chat_id, cur_title, st.session_state.messages)
                    loaded = load_chat_from_db(chat_id)
                    if loaded is not None:
                        st.session_state.messages = loaded
                        st.session_state.current_chat_id = chat_id
                        st.session_state.documents = []
                        st.session_state.processed_file_names = set()
                        st.rerun()
            with col2:
                if st.button("✕", key=f"del_{chat_id}"):
                    delete_chat_from_db(chat_id)
                    if is_current:
                        st.session_state.messages = []
                        st.session_state.current_chat_id = str(uuid.uuid4())
                    st.rerun()

    with st.expander("🔌  Provider & model", expanded=True):
        provider = st.selectbox(
            "AI Provider",
            ["Ollama (local)", "Gemini", "OpenAI", "Anthropic", "Groq"],
            label_visibility="collapsed",
        )

        api_key = None
        model = None

        if provider == "Ollama (local)":
            available_models = get_ollama_models()
            if available_models:
                model = st.selectbox("Model", available_models)
            else:
                st.warning("⚠️ Couldn't reach Ollama. Is it running? (`ollama serve`)")
                model = st.text_input("Model name", value="llama3.2")

            if model and not is_ollama_vision_model(model):
                st.caption(
                    f"ℹ️ '{model}' doesn't look like a vision model — images you upload "
                    "won't be understood unless you switch to a model like `llava` "
                    "or `llama3.2-vision` (run `ollama pull llava` to get one)."
                )

        elif provider == "Gemini":
            model = st.selectbox("Model", GEMINI_MODELS)
            api_key = st.text_input("Gemini API Key", type="password")

        elif provider == "OpenAI":
            model = st.selectbox("Model", OPENAI_MODELS)
            api_key = st.text_input("OpenAI API Key", type="password")

        elif provider == "Anthropic":
            model = st.selectbox("Model", ANTHROPIC_MODELS)
            api_key = st.text_input("Anthropic API Key", type="password")

        elif provider == "Groq":
            model = st.selectbox("Model", GROQ_MODELS)
            api_key = st.text_input("Groq API Key", type="password")
            if model and not is_groq_vision_model(model):
                st.caption(
                    f"ℹ️ '{model}' doesn't look like a vision model — images you upload "
                    "won't be understood unless you switch to a model like "
                    "`meta-llama/llama-4-scout-17b-16e-instruct`."
                )

        # Active provider chip preview
        style = PROVIDER_STYLE[provider]
        st.markdown(
            f'<div style="display:inline-flex; align-items:center; gap:6px; '
            f'background-color:{style["bg"]}; border:1px solid {style["border"]}; '
            f'border-radius:14px; padding:5px 12px; margin-top:6px; font-size:12px; '
            f'color:{style["color"]};">{style["icon"]} {style["label"]} · {model}</div>',
            unsafe_allow_html=True,
        )

    with st.expander("🔎  Web search", expanded=False):
        st.caption(
            "Auto-detects questions needing current info (e.g. 'latest', 'today', "
            "'current', 'news') and tries a live search before answering. "
            "Best-effort — DuckDuckGo may occasionally block or rate-limit it."
        )

    with st.expander("📄  Documents", expanded=bool(st.session_state.documents)):
        st.caption("Upload PDF, Word, PowerPoint, Excel, images, text, Markdown, or Python files.")
        uploaded_files = st.file_uploader(
            "Upload files",
            type=["pdf", "docx", "pptx", "xlsx", "txt", "md", "py", "png", "jpg", "jpeg"],
            accept_multiple_files=True,
            label_visibility="collapsed",
        )

        if uploaded_files:
            for f in uploaded_files:
                if f.name in st.session_state.processed_file_names:
                    continue  # already processed this exact upload, skip on rerun
                if len(st.session_state.documents) >= MAX_DOCS_ACTIVE:
                    st.warning(f"⚠️ Limit of {MAX_DOCS_ACTIVE} active documents reached — remove one first.")
                    break

                with st.spinner(f"Processing {f.name}..."):
                    doc = process_uploaded_file(f)
                st.session_state.documents.append(doc)
                st.session_state.processed_file_names.add(f.name)

                if doc["image_render_attempted"] and doc["image_render_failed"]:
                    st.warning(
                        f"⚠️ '{f.name}': extracted text okay, but couldn't render slide/sheet "
                        "images (LibreOffice not found or conversion failed). Install LibreOffice "
                        "and ensure `soffice` is on your PATH for visual rendering."
                    )
                else:
                    st.success(f"Loaded: {f.name}")

        if st.session_state.documents:
            st.markdown("**Active documents:**")
            for i, doc in enumerate(st.session_state.documents):
                img_note = f" · {len(doc['images'])} image(s)" if doc["images"] else ""
                col1, col2 = st.columns([5, 1])
                with col1:
                    st.caption(f"📎 {doc['name']}{img_note}")
                with col2:
                    if st.button("✕", key=f"rmdoc_{i}"):
                        removed = st.session_state.documents.pop(i)
                        st.session_state.processed_file_names.discard(removed["name"])
                        st.rerun()

            if st.button("Remove all documents", use_container_width=True):
                st.session_state.documents = []
                st.session_state.processed_file_names = set()
                st.rerun()


# =============================================================================
# MAIN CHAT UI
# =============================================================================
st.markdown("### General chatbot")

# Display existing chat history using custom bubbles
for msg in st.session_state.messages:
    if msg["role"] == "user":
        render_bubble("user", msg["content"], msg.get("timestamp", ""))
    else:
        provider_style = PROVIDER_STYLE.get(msg.get("provider", "Anthropic"), PROVIDER_STYLE["Anthropic"])
        render_bubble(
            "assistant",
            msg["content"],
            msg.get("timestamp", ""),
            provider_info=provider_style,
            search_used=msg.get("search_used"),
        )

# Chat input
user_input = st.chat_input("Type your message...")

if user_input:
    now = datetime.now().strftime("%I:%M %p")

    # Show user message immediately
    render_bubble("user", user_input, now)
    st.session_state.messages.append({"role": "user", "content": user_input, "timestamp": now})

    # Build the message list to send to the provider. Strip UI-only fields
    # (timestamp, provider, search_used) since those aren't part of the chat protocol.
    messages_to_send = [
        {"role": m["role"], "content": m["content"]} for m in st.session_state.messages
    ]

    # ---- Combine all active documents into one shared context ----
    # Per the user's choice, every uploaded document stays active and is
    # considered together on every question (no per-question selection).
    all_images = []
    vision_capable = (
        provider in ("Gemini", "OpenAI", "Anthropic")
        or (provider == "Ollama (local)" and is_ollama_vision_model(model or ""))
        or (provider == "Groq" and is_groq_vision_model(model or ""))
    )
    image_warning = None

    if st.session_state.documents:
        text_sections = []
        for doc in st.session_state.documents:
            if doc["text"]:
                text_sections.append(f"=== File: {doc['name']} ===\n{doc['text']}")
            if doc["images"]:
                if vision_capable:
                    all_images.extend(doc["images"])
                else:
                    image_warning = (
                        f"Note: '{model}' on {provider} can't view images, so visual "
                        f"content from {doc['name']} (and any other image/slide files) "
                        f"was skipped — only extracted text was used."
                    )

        if text_sections:
            doc_context = (
                f"The user has uploaded {len(st.session_state.documents)} document(s). "
                f"Use their content below to answer questions when relevant. "
                f"If the question is unrelated to the documents, just answer normally.\n\n"
                + "\n\n".join(text_sections)
                + f"\n\nUser question: {user_input}"
            )
            messages_to_send[-1]["content"] = doc_context
        elif all_images:
            # Images only, no extracted text (e.g. a single image upload)
            messages_to_send[-1]["content"] = user_input

        if all_images:
            messages_to_send[-1]["images"] = all_images

    # ---- Auto-detected web search ----
    search_used = None  # None = not attempted, True = succeeded, False = attempted but failed
    if needs_web_search(user_input):
        search_used = False
        with st.spinner("🔎 Searching the web for current info..."):
            search_results = web_search_duckduckgo(user_input)

        if search_results:
            search_used = True
            search_context = (
                f"Live web search results for the query '{user_input}' "
                f"(retrieved just now):\n\n{search_results}\n\n"
                f"Use these results to answer the user's question with current "
                f"information. If the results don't fully answer it, say so and "
                f"answer with your best general knowledge instead.\n\n"
                f"User question: {user_input}"
            )
            current_content = messages_to_send[-1]["content"]
            if current_content != user_input:
                messages_to_send[-1]["content"] = current_content + "\n\n" + search_context
            else:
                messages_to_send[-1]["content"] = search_context

    if image_warning:
        st.caption(f"ℹ️ {image_warning}")

    # Call the selected provider
    with st.spinner("Thinking..."):
        if provider == "Ollama (local)":
            reply = call_ollama(messages_to_send, model)
        elif provider == "Gemini":
            reply = call_gemini(messages_to_send, model, api_key)
        elif provider == "OpenAI":
            reply = call_openai(messages_to_send, model, api_key)
        elif provider == "Anthropic":
            reply = call_anthropic(messages_to_send, model, api_key)
        elif provider == "Groq":
            reply = call_groq(messages_to_send, model, api_key)
        else:
            reply = "⚠️ Unknown provider selected."

    reply_time = datetime.now().strftime("%I:%M %p")
    render_bubble(
        "assistant",
        reply,
        reply_time,
        provider_info=PROVIDER_STYLE[provider],
        search_used=search_used,
    )

    st.session_state.messages.append({
        "role": "assistant",
        "content": reply,
        "timestamp": reply_time,
        "provider": provider,
        "search_used": search_used,
    })

    # ---- Persist this chat to the database ----
    # Saved/updated after every exchange so nothing is lost if the app closes.
    chat_title = generate_chat_title(st.session_state.messages[0]["content"])
    save_chat_to_db(st.session_state.current_chat_id, chat_title, st.session_state.messages)


